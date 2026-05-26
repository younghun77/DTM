"""Generate the DTM Factory Kit overview PowerPoint deck.

Run:
    py -3 tools/make_pptx.py
Output:
    dist/DTM_Factory_Kit_Overview.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.abspath(os.path.join(HERE, "..", "dist", "DTM_Factory_Kit_Overview.pptx"))

NAVY = RGBColor(0x16, 0x3A, 0x6A)
TEAL = RGBColor(0x1F, 0x8F, 0x5D)
RED = RGBColor(0xCC, 0x33, 0x33)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xEE, 0xEE, 0xEE)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Color band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  prs.slide_width, Inches(2.0))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6),
                                  prs.slide_width - Inches(1.2), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sb = slide.shapes.add_textbox(Inches(0.6), Inches(2.2),
                                  prs.slide_width - Inches(1.2), Inches(1.0))
    p2 = sb.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.runs[0].font.size = Pt(20)
    p2.runs[0].font.color.rgb = NAVY

    foot = slide.shapes.add_textbox(Inches(0.6), prs.slide_height - Inches(0.8),
                                    prs.slide_width - Inches(1.2), Inches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "nRF52840 Dongle | Bluetooth Direct Test Mode | Standalone Distribution"
    fp.runs[0].font.size = Pt(12)
    fp.runs[0].font.color.rgb = GREY


def add_section_slide(prs: Presentation, title: str, bullets: list[str],
                      accent: RGBColor = NAVY) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15),
                                  prs.slide_width - Inches(0.8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(26); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1),
                                    prs.slide_width - Inches(1.0),
                                    prs.slide_height - Inches(1.4))
    tf = body.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # indent level: lines starting with "- " become sub-bullets
        if bullet.startswith("- "):
            para.level = 1
            para.text = bullet[2:]
        else:
            para.level = 0
            para.text = bullet
        for r in para.runs:
            r.font.size = Pt(16 if para.level == 0 else 13)
            r.font.color.rgb = NAVY if para.level == 0 else GREY


def add_table_slide(prs: Presentation, title: str,
                    headers: list[str], rows: list[list[str]],
                    accent: RGBColor = NAVY) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15),
                                  prs.slide_width - Inches(0.8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.runs[0].font.size = Pt(26); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    nrows = len(rows) + 1
    ncols = len(headers)
    table_shape = slide.shapes.add_table(
        nrows, ncols,
        Inches(0.4), Inches(1.1),
        prs.slide_width - Inches(0.8),
        min(Inches(0.45) * nrows, prs.slide_height - Inches(1.4)))
    table = table_shape.table
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.text = h
        for r in c.text_frame.paragraphs[0].runs:
            r.font.bold = True; r.font.size = Pt(13)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.fill.solid(); c.fill.fore_color.rgb = accent
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j); c.text = val
            for r in c.text_frame.paragraphs[0].runs:
                r.font.size = Pt(11); r.font.color.rgb = NAVY


def build() -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs,
                    "DTM Factory Kit",
                    "nRF52840 Dongle 기반 Direct Test Mode 자동 RX 검사 패키지")

    # 2. Background & Goal
    add_section_slide(prs, "1. 배경 & 목적", [
        "기존: nRF Connect Desktop의 Direct Test Mode 앱 + 별도 SDK 설치 필요",
        "- 작업자 PC마다 환경 구축이 필요해 라인 투입 지연",
        "- DUT TX 제어와 RX 측정이 분리되어 수동 운용",
        "목표: 별도 SDK 없이 nRF Dongle만 꽂으면 동작하는 단독 패키지",
        "- 자동 RX 측정 + DUT 제어(SSH/Ethernet) 통합 GUI",
        "- 실패 자동 감지 / 로그 수집 / 결과 그래프화",
        "- Standalone(외부 TX 장비) 모드 지원",
    ])

    # 3. System Architecture
    add_section_slide(prs, "2. 시스템 구조", [
        "PC (Windows)",
        "- run_gui.bat -> tools\\dtm_rx_runner.py (Tkinter GUI)",
        "- dut_control.py (paramiko SSH + TCP 20000 reboot 프레임)",
        "- bin\\nrfutil.exe (DFU 플래시)",
        "Dongle (nRF52840)",
        "- DTM 펌웨어 + Legacy USB CDC-ACM 'Nordic DTM USB'",
        "- 19200 8N1, Bluetooth 2-wire DTM 프로토콜",
        "DUT (Test Sample)",
        "- SSH: bt_tx_test_39ch.sh / bt_test_off.sh",
        "- TCP 20000: KU 바이너리 reboot 프레임",
        "- /var/data/btman/bt_test.log, bt_bootstrap.log",
    ], accent=TEAL)

    # 4. Distribution Layout
    add_section_slide(prs, "3. 배포 패키지 구성", [
        "dtm_factory_kit/",
        "- flash_dongle.bat            nRF52840 Dongle DFU 플래시",
        "- run_gui.bat                 RX 테스트 GUI 실행 (pip 자동 설치)",
        "- README.md, CHANGELOG.md     사용·이력 문서",
        "- bin/nrfutil.exe             Nordic 통합 도구 (내장)",
        "- firmware/dtm_dongle.zip     dongle DFU 패키지",
        "- tools/dtm_rx_runner.py      통합 GUI",
        "- tools/dut_control.py        DUT 제어 라이브러리 (+ CLI)",
        "- tools/private_key.pem       OpenSSH 키 (선택)",
        "한 폴더만 압축 해제 → 즉시 사용 (Python 3.8+ 외 추가 설치 불필요)",
    ])

    # 5. GUI features
    add_section_slide(prs, "4. GUI 주요 기능", [
        "메인 버튼: START RX TEST / END RX TEST / REBOOT DUT / DTM Reset / Open CSV",
        "AUTO 행: Iterations, RX duration, Cooldown, AUTO RUN, AUTO RUN (no reboot), STOP, Plot CSV",
        "Mode Switch: DUT-LINK <-> STANDALONE (dongle only) 토글",
        "- Standalone 시 AUTO/REBOOT 버튼 자동 숨김",
        "결과: D:\\factory\\YY-MM-DD\\rx_result.csv 자동 저장",
        "그래프: matplotlib 있으면 라인+평균선, 없으면 Tk 캔버스 폴백",
    ], accent=TEAL)

    # 6. Auto failure detection
    add_section_slide(prs, "5. 자동 실패 감지 & 로그 수집", [
        "AUTO 루프 중 다음 조건에서 즉시 중단 + 아티팩트 수집",
        "- rx_count == 0  (DUT TX 미발생 / FW crash 의심)",
        "- bt_test.log 에서 [ERROR][BT] 검출",
        "- bt_bootstrap.log 에서 'Fail, total boot time' 검출",
        "- 로그에서 firmware-crash 시그니처 ff fd 01 08 55 00 검출",
        "수집 절차:",
        "- SFTP로 /var/data/btman/bt_test.log, bt_bootstrap.log 다운로드",
        "- analyze_dut_logs() 가 매칭 라인을 quote 하여 출력",
        "- SSH cat /proc/tty/driver/serial 결과로 serial 상태 확인",
        "- D:\\factory\\YY-MM-DD\\logs\\HHMMSS_iterN_<reason>\\FAILURE_SUMMARY.txt 작성",
    ], accent=RED)

    # 7. Reduction & effect (table)
    add_table_slide(prs, "6. 절감 내용 & 효과",
                    headers=["항목", "기존", "DTM Factory Kit", "효과"],
                    rows=[
        ["환경 셋업", "nRF Connect Desktop + SDK 설치 (~30 min/대)", "zip 해제 후 실행 (<2 min)", "초기 셋업 -90% 이상"],
        ["RX 1회 측정", "수동 (DUT TX 시작 → 앱에서 RX → 결과 기록)", "버튼 1회 (START/END)", "조작 단계 5→1, 휴먼 에러 감소"],
        ["반복 1000회 실행", "수동 반복 / 자체 스크립트 필요", "AUTO RUN + STOP, CSV 자동 적재", "야간 무인 평가 가능"],
        ["실패 분석", "재현 → 수동 로그 수집 (10~30 min)", "자동 정지 + 로그 다운로드/분석", "분석 시간 -80%"],
        ["외부 TX 장비 연동", "별도 툴 필요", "Standalone 스위치로 즉시 전환", "검사 환경 유연성 향상"],
        ["결과 시각화", "Excel 수작업", "Plot CSV 버튼 (라인+평균)", "리뷰 즉시성"],
    ])

    # 8. Workflow comparison
    add_section_slide(prs, "7. 작업 흐름 비교 (Before / After)", [
        "BEFORE",
        "- 1) SDK / nRF Connect 설치  2) DTM 앱 실행  3) DUT TX 수동  4) 채널 별 측정",
        "- 5) 결과를 별도 시트에 기록  6) 실패 시 수동 SSH 로 로그 수집",
        "AFTER",
        "- 1) flash_dongle.bat (최초 1회)  2) run_gui.bat",
        "- 3) AUTO RUN 또는 START/END 버튼  4) CSV/그래프 자동",
        "- 5) 실패 자동 정지 + bt_test.log/bt_bootstrap.log/serial 덤프 자동 수집",
    ], accent=TEAL)

    # 9. Key technical highlights
    add_section_slide(prs, "8. 기술적 포인트", [
        "Persistent SSHSession 도입으로 매 명령마다 핸드셰이크 제거",
        "reboot_dut(): reboot 이후 TimeoutError를 정상으로 처리 → 안정적 무한 루프",
        "fire-and-forget bt_tx_test_39ch.sh 실행 + 3s startup_delay → RX 윈도우 보장",
        "fetch_dut_logs (SFTP) + analyze_dut_logs (텍스트/바이너리 패턴) 일체화",
        "Plot CSV: matplotlib 또는 순수 Tk Canvas 폴백 → 클린 PC에서도 그래프 가능",
        "Standalone 모드: SSH/이더넷 호출 전부 no-op 처리, AUTO 버튼 동적 숨김",
    ])

    # 10. Roadmap & contact
    add_section_slide(prs, "9. 향후 개선 후보", [
        "bt_bootstrap.log STOP 문자열 감지 자동 정지 옵션 토글화",
        "채널별 RX 결과 분리 / 멀티 그래프",
        "DUT IP, SSH 사용자, 기본 채널을 config.ini로 외부화",
        "GUI 로그를 파일로도 동시 저장 (logs/YYYY-MM-DD.log)",
        "GitHub Release 태깅 자동화 (gh release create)",
    ])

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
