"""Generate the DTM RX Runner presentation deck.

Focus:
  * Test Sync between Test Sample (DUT) and Dongle RX
  * Script selection per OEM sample (MANUAL mode)
  * Architecture diagram (TestSample <-> RX Runner <-> Dongle)
  * Feature list & explicit Test Sequence

Run:
    py -3 tools/make_rx_runner_pptx.py
Output:
    dist/DTM_RX_Runner.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.abspath(os.path.join(HERE, "..", "dist", "DTM_RX_Runner.pptx"))

# ---- Palette (matches the GUI's "slate" look) -------------------------------
NAVY = RGBColor(0x0F, 0x17, 0x2A)   # bg
SLATE = RGBColor(0x1E, 0x29, 0x3B)  # panel
SLATE2 = RGBColor(0x33, 0x41, 0x55)
SKY = RGBColor(0x38, 0xBD, 0xF8)    # accent
OK = RGBColor(0x22, 0xC5, 0x5E)
WARN = RGBColor(0xF5, 0x9E, 0x0B)
ERR = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x94, 0xA3, 0xB8)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)


# ---------- helpers ----------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, x, y, w, h, text, *, size=14, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def _title_bar(slide, prs, title, accent=SKY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, Inches(0.9))
    _fill(bar, SLATE)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85),
                                    prs.slide_width, Inches(0.05))
    _fill(stripe, accent)
    _textbox(slide, Inches(0.4), Inches(0.18),
             prs.slide_width - Inches(0.8), Inches(0.6),
             title, size=24, bold=True, color=WHITE)


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    _fill(bg, NAVY)
    return slide


def _bullets(slide, x, y, w, h, items, *, size=15, color=TEXT):
    """items: list[str].  '- ' prefix becomes a sub-bullet."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("- "):
            para.level = 1
            para.text = "• " + line[2:]
            for r in para.runs:
                r.font.size = Pt(size - 2)
                r.font.color.rgb = GREY
                r.font.name = "Segoe UI"
        else:
            para.level = 0
            para.text = "▸ " + line
            for r in para.runs:
                r.font.size = Pt(size)
                r.font.color.rgb = color
                r.font.bold = True
                r.font.name = "Segoe UI"
        para.space_after = Pt(4)
    return tb


def _card(slide, x, y, w, h, title, body_lines, accent=SKY):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _fill(bg, SLATE)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, w, Inches(0.38))
    _fill(head, accent)
    _textbox(slide, x + Inches(0.15), y + Inches(0.04),
             w - Inches(0.3), Inches(0.32),
             title, size=13, bold=True, color=NAVY)
    _bullets(slide, x + Inches(0.15), y + Inches(0.5),
             w - Inches(0.3), h - Inches(0.6),
             body_lines, size=12)


def _node(slide, x, y, w, h, title, sub, *, color=SKY, fg=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _fill(box, color)
    _textbox(slide, x, y + Inches(0.08), w, Inches(0.4),
             title, size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)
    _textbox(slide, x, y + Inches(0.55), w, h - Inches(0.6),
             sub, size=10, color=fg, align=PP_ALIGN.CENTER)


def _arrow(slide, x1, y1, x2, y2, *, color=SKY, label=None,
           label_dy=Inches(-0.28)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(2.25)
    # arrow head
    line.line._get_or_add_ln().set("cap", "flat")
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    if label:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        _textbox(slide, mx - Inches(1.4), my + label_dy,
                 Inches(2.8), Inches(0.32),
                 label, size=10, bold=True,
                 color=color, align=PP_ALIGN.CENTER)


# ---------- slide builders ---------------------------------------------------
def slide_title(prs):
    slide = _blank_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(3.0))
    _fill(band, SLATE)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    0, Inches(2.95), prs.slide_width,
                                    Inches(0.08))
    _fill(stripe, SKY)

    _textbox(slide, Inches(0.6), Inches(0.7),
             prs.slide_width - Inches(1.2), Inches(1.2),
             "DTM RX Runner", size=44, bold=True, color=WHITE)
    _textbox(slide, Inches(0.6), Inches(1.7),
             prs.slide_width - Inches(1.2), Inches(0.6),
             "Bluetooth Direct Test Mode  —  Dongle ↔ Test Sample 통합 RX 자동화",
             size=20, color=SKY)
    _textbox(slide, Inches(0.6), Inches(2.3),
             prs.slide_width - Inches(1.2), Inches(0.5),
             "Test Sync · Script Selection · Auto Loop · Failure Capture",
             size=14, color=GREY)

    _textbox(slide, Inches(0.6), Inches(6.6),
             prs.slide_width - Inches(1.2), Inches(0.4),
             "nRF52840 Dongle  |  Python/Tk GUI  |  paramiko SSH + KU TCP control",
             size=11, color=GREY)


def slide_why(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "1. 왜 만들었나 (Why)", accent=SKY)

    # Comparison table-like cards
    col_w = (prs.slide_width - Inches(1.2)) // 2
    y = Inches(1.2)
    h = Inches(5.6)
    _card(slide, Inches(0.4), y, col_w, h,
          "기존 도구 (nRF Connect DTM)",
          [
              "Dongle 한쪽만 RX/TX 제어",
              "- DUT(Test Sample)는 운영자가 수동 트리거",
              "- DUT TX 시작과 RX 윈도우 사이 timing jitter",
              "단발 측정 위주",
              "- 결과는 운영자가 따로 기록",
              "OEM 별 스크립트가 다르면 코드 수정 필요",
              "실패 원인 분석은 수동 SSH/로그 다운로드",
          ], accent=GREY)
    _card(slide, Inches(0.8) + col_w, y, col_w, h,
          "DTM RX Runner",
          [
              "Dongle RX + DUT TX 동시 제어 (Test Sync)",
              "- SSH로 DUT 스크립트를 트리거 → RX 윈도와 자동 정렬",
              "- 측정 시점이 항상 같아 재현성 확보",
              "AUTO RUN N회 반복 + CSV/그래프 자동 저장",
              "OEM/샘플마다 Script Selection (UI)",
              "- 코드 수정 없이 신규 샘플 즉시 측정 가능",
              "실패 자동 감지 → 로그 SFTP + Summary 작성",
          ], accent=SKY)


def slide_features(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "2. 주요 특징 (Features)", accent=SKY)

    cards = [
        ("Test Sync (테스트 동기화)", OK, [
            "START RX 한 번 → DUT TX + Dongle RX 동시 시작",
            "- DUT TX 스크립트 SSH 트리거 → startup delay → Receiver Test",
            "- AUTO RUN 반복 시에도 매 iteration 동기 보장",
            "- 운영자 손동작에 의한 timing jitter 제거",
        ]),
        ("Script Selection (스크립트 선택)", SKY, [
            "MANUAL 모드에서 DUT의 *.sh 목록 자동 fetch",
            "- TX script / Test End script를 GUI에서 클릭 선택",
            "- 신규 OEM 샘플 추가 시 코드 수정 불필요",
            "- 'Select script' 버튼 + Combobox로 동기 노출",
        ]),
        ("3가지 동작 모드", VIOLET, [
            "DUT-LINK  : 표준 BMW Telematics 설정",
            "MANUAL    : SSH/Ethernet/Reboot frame/스크립트 사용자 지정",
            "STANDALONE: Dongle 단독 (외부 신호원 + RX 측정)",
        ]),
        ("AUTO Loop & 실패 자동 진단", WARN, [
            "Iterations / RX dur / Cooldown 지정 후 무인 반복",
            "- rx_count==0 또는 known pattern 검출 시 자동 중단",
            "- bt_test.log / bt_bootstrap.log SFTP 다운로드",
            "- FAILURE_SUMMARY.txt 자동 생성",
        ]),
        ("운영자 편의성", BLUE, [
            "Dongle 자동 탐지(VID 0x1915), 실패 시 수동 picker",
            "STAGE 인디케이터 + 활성 버튼 highlight ring",
            "CSV 자동 저장 (results/YY-MM-DD/rx_result.csv)",
            "Plot CSV: matplotlib / Tk Canvas fallback",
        ]),
        ("안정성 / 확장성", ERR, [
            "Persistent SSH session (handshake 1회) → jitter↓",
            "Reboot frame 유효성 트레이스 → 잘못된 조작 차단",
            "신규 reboot 프로토콜은 hex frame만 교체",
            "환경변수 DTM_RX_RESULT_BASE로 결과 경로 오버라이드",
        ]),
    ]
    # 3x2 grid
    cw = (prs.slide_width - Inches(1.6)) // 3
    ch = Inches(2.7)
    x0 = Inches(0.4)
    y0 = Inches(1.15)
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)
    for idx, (title, color, lines) in enumerate(cards):
        col = idx % 3
        row = idx // 3
        x = x0 + (cw + gap_x) * col
        y = y0 + (ch + gap_y) * row
        _card(slide, x, y, cw, ch, title, lines, accent=color)


def slide_architecture(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "3. 아키텍처 (TestSample ↔ RX Runner ↔ Dongle)",
               accent=SKY)

    # ----- PC frame -----
    pc_x = Inches(0.6); pc_y = Inches(1.2)
    pc_w = Inches(5.4); pc_h = Inches(4.0)
    pc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                pc_x, pc_y, pc_w, pc_h)
    _fill(pc, SLATE)
    _textbox(slide, pc_x, pc_y + Inches(0.05), pc_w, Inches(0.35),
             "Operator PC (Windows)", size=12, bold=True,
             color=SKY, align=PP_ALIGN.CENTER)

    # GUI box
    gui_x = pc_x + Inches(0.25); gui_y = pc_y + Inches(0.55)
    gui_w = pc_w - Inches(0.5); gui_h = Inches(3.2)
    gui = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 gui_x, gui_y, gui_w, gui_h)
    _fill(gui, SLATE2)
    _textbox(slide, gui_x, gui_y + Inches(0.05), gui_w, Inches(0.35),
             "DTM RX Runner  (tools/dtm_rx_runner.py)",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Two engine boxes inside GUI
    eng_w = (gui_w - Inches(0.45)) // 2
    eng_h = Inches(1.0)
    eng_y = gui_y + Inches(0.55)
    _node(slide, gui_x + Inches(0.15), eng_y, eng_w, eng_h,
          "DTM Engine", "2-wire encode/decode\nCMD_RESET / RX / END",
          color=SKY)
    _node(slide, gui_x + Inches(0.3) + eng_w, eng_y, eng_w, eng_h,
          "DUT Control", "paramiko SSH session\nKU TCP reboot frame",
          color=OK)

    # Sequencer box
    seq_y = eng_y + eng_h + Inches(0.2)
    _node(slide, gui_x + Inches(0.15), seq_y,
          gui_w - Inches(0.3), Inches(0.95),
          "Sequencer / AUTO Loop",
          "Test Sync · Script Selection · CSV · Failure Capture",
          color=VIOLET, fg=WHITE)

    # ----- Dongle -----
    d_x = Inches(7.6); d_y = Inches(1.3)
    d_w = Inches(2.5); d_h = Inches(2.4)
    _node(slide, d_x, d_y, d_w, d_h,
          "nRF52840 Dongle",
          "DTM firmware\nUSB CDC-ACM\n19200 8N1",
          color=SKY)

    # ----- DUT -----
    t_x = Inches(10.3); t_y = Inches(3.8)
    t_w = Inches(2.6); t_h = Inches(2.6)
    _node(slide, t_x, t_y, t_w, t_h,
          "Test Sample (DUT)",
          "BT chipset + Linux\n/opt/factory/.../*.sh\nbt_tx_test_*.sh, bt_test_off.sh",
          color=OK)

    # ----- Arrows -----
    # USB CDC-ACM: PC GUI -> Dongle
    _arrow(slide,
           gui_x + gui_w, eng_y + eng_h // 2,
           d_x, d_y + d_h // 2,
           color=SKY, label="USB CDC-ACM  (DTM 2-wire)")

    # SSH: PC GUI -> DUT
    _arrow(slide,
           gui_x + gui_w, seq_y + Inches(0.3),
           t_x, t_y + Inches(0.7),
           color=OK,
           label="SSH (key)  TX/Off script · SFTP logs")

    # TCP control: PC GUI -> DUT
    _arrow(slide,
           gui_x + gui_w, seq_y + Inches(0.7),
           t_x, t_y + Inches(1.6),
           color=WARN,
           label="TCP :20000  KU reboot frame")

    # BLE: Dongle <-> DUT
    _arrow(slide,
           d_x + d_w // 2, d_y + d_h,
           t_x + t_w // 2, t_y,
           color=VIOLET,
           label="2.4 GHz BLE  (TX ↔ RX)")

    # Legend
    _textbox(slide, Inches(0.6), Inches(6.5),
             prs.slide_width - Inches(1.2), Inches(0.4),
             "측정 RF 경로 = Dongle ↔ Test Sample (BLE) /  제어 경로 = PC → Dongle(USB), PC → DUT(SSH+TCP)",
             size=11, color=GREY, align=PP_ALIGN.CENTER)


def slide_sequence_single(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "4. Test Sequence — 단일 측정 (START RX → END RX)",
               accent=OK)

    # 4 lanes
    lane_y = Inches(1.2)
    lane_h = Inches(0.5)
    headers = [("Operator", SKY), ("DTM RX Runner", VIOLET),
               ("Dongle", BLUE), ("Test Sample (DUT)", OK)]
    lane_w = (prs.slide_width - Inches(0.8)) // 4
    for i, (name, color) in enumerate(headers):
        x = Inches(0.4) + lane_w * i
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, lane_y, lane_w, lane_h)
        _fill(head, color)
        _textbox(slide, x, lane_y + Inches(0.08), lane_w, Inches(0.4),
                 name, size=13, bold=True, color=NAVY if color != VIOLET else WHITE,
                 align=PP_ALIGN.CENTER)
        # vertical lane line
        col = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + lane_w // 2 - Emu(9525), lane_y + lane_h,
                                     Emu(19050), Inches(5.2))
        _fill(col, SLATE2)

    # Steps as horizontal arrows + text
    def col_center(i):
        return Inches(0.4) + lane_w * i + lane_w // 2

    steps = [
        # (from_lane, to_lane, text, color, y_offset_in)
        (0, 1, "① START RX 클릭", SKY, 0.0),
        (1, 3, "② SSH exec: TX script (async)", OK, 0.6),
        (None, None, "⏱ startup_delay 3s (DUT TX 안정화)", GREY, 1.2),
        (1, 2, "③ DTM CMD_RECEIVER (ch, len, PRBS9)", BLUE, 1.5),
        (None, None, "RX 윈도 동안 BLE 패킷 수신 (Dongle ⇐ DUT)", VIOLET, 2.1),
        (0, 1, "④ END RX 클릭", SKY, 2.7),
        (1, 2, "⑤ DTM CMD_END → packet_count 회수", BLUE, 3.3),
        (None, None, "⑥ CSV append  (results/YY-MM-DD/rx_result.csv)", OK, 3.9),
        (1, 3, "⑦ SSH exec: Test End script (TX OFF)", OK, 4.5),
    ]
    y0 = lane_y + lane_h + Inches(0.2)
    for s in steps:
        from_i, to_i, label, color, dy = s
        y = y0 + Inches(dy)
        if from_i is None:
            # full-width note
            _textbox(slide, Inches(0.6), y,
                     prs.slide_width - Inches(1.2), Inches(0.35),
                     "    " + label, size=12, bold=True,
                     color=color, align=PP_ALIGN.CENTER)
        else:
            x1 = col_center(from_i); x2 = col_center(to_i)
            _arrow(slide, x1, y + Inches(0.18), x2, y + Inches(0.18),
                   color=color, label=label, label_dy=Inches(-0.28))


def slide_sequence_auto(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "5. Test Sequence — AUTO RUN (N회 반복, with reboot)",
               accent=BLUE)

    # Pseudo-code box
    code_x = Inches(0.5); code_y = Inches(1.2)
    code_w = prs.slide_width - Inches(1.0); code_h = Inches(5.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 code_x, code_y, code_w, code_h)
    _fill(box, SLATE)

    lines = [
        ("for i in 1..N:", SKY, True),
        ("    [DUT TX START]   SSH: TX script 실행 (async, blocking 아님)", TEXT, False),
        ("    [SETTLE]         3s startup delay  ─ DUT가 실제 송신 시작할 시간", GREY, False),
        ("    [DONGLE RESET]   CMD_RESET", TEXT, False),
        ("    [DONGLE RX ON]   CMD_RECEIVER(ch, len, PRBS9)", BLUE, True),
        ("    [RX WINDOW]      RX duration 동안 수신   (STOP 버튼으로 중단 가능)", VIOLET, False),
        ("    [DONGLE END]     CMD_END → packet_count 회수", BLUE, True),
        ("    [CSV APPEND]     results/YY-MM-DD/rx_result.csv", OK, True),
        ("", TEXT, False),
        ("    if rx_count == 0:                       # 즉시 실패 처리", ERR, True),
        ("        fetch_dut_logs(SFTP) → analyze → FAILURE_SUMMARY.txt", ERR, False),
        ("        break", ERR, False),
        ("", TEXT, False),
        ("    if reboot_between:                       # AUTO RUN (with reboot)", WARN, True),
        ("        SSH: Test End script (HCI 해제) → 1s settle", WARN, False),
        ("        TCP: SERVICE_ENABLE frame → ACK", WARN, False),
        ("        TCP: REBOOT frame          → DUT reboot", WARN, False),
        ("    else:                                    # AUTO (no reboot)", VIOLET, True),
        ("        SSH: Test End script", VIOLET, False),
        ("", TEXT, False),
        ("    [LOG SCAN]       known failure pattern? (SSH 불안정 시 3회 retry)", GREY, False),
        ("    [COOLDOWN]       N초 대기 (STOP 가능)", GREY, False),
        ("", TEXT, False),
        ("[FINISH]            STAGE: idle, AUTO 버튼 복구", SKY, True),
    ]
    tb = slide.shapes.add_textbox(code_x + Inches(0.25), code_y + Inches(0.2),
                                  code_w - Inches(0.5), code_h - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        for r in p.runs:
            r.font.name = "Cascadia Mono"
            r.font.size = Pt(12)
            r.font.bold = bold
            r.font.color.rgb = color
        p.space_after = Pt(0)


def slide_manual_setup(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "6. MANUAL 모드 셋업 시퀀스 (신규 OEM 샘플 대응)",
               accent=VIOLET)

    steps = [
        ("①", "Mode 토글",
         "DUT-LINK → MANUAL  (TEST CONFIG 카드의 Mode 버튼)"),
        ("②", "DUT settings… 다이얼로그 오픈",
         "SSH host / user / PEM key (Browse…)\nEthernet host:port  (reboot 대상)\nReboot frame (hex)  ← 비우면 REBOOT/AUTO RUN 자동 비활성화\nScript dir"),
        ("③", "Connect SSH",
         "paramiko로 연결 → 원격 디렉터리의 *.sh 자동 listing"),
        ("④", "Select script",
         "TX script / Test End script 클릭 선택\n→ 코드 수정 없이 OEM별 스크립트 적용"),
        ("⑤", "Apply settings",
         "활성 설정 로그 출력, STAGE에 'MANUAL applied' 표시"),
        ("⑥", "Close → 측정 시작",
         "START RX / AUTO RUN 등 평상시와 동일하게 사용"),
    ]
    y = Inches(1.2)
    h = Inches(0.85)
    gap = Inches(0.08)
    for idx, (num, title, body) in enumerate(steps):
        yy = y + (h + gap) * idx
        # number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(0.5), yy + Inches(0.1),
                                      Inches(0.65), Inches(0.65))
        _fill(circ, VIOLET)
        _textbox(slide, Inches(0.5), yy + Inches(0.18),
                 Inches(0.65), Inches(0.5),
                 num, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # content card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.3), yy,
                                      prs.slide_width - Inches(1.8), h)
        _fill(card, SLATE)
        _textbox(slide, Inches(1.5), yy + Inches(0.08),
                 prs.slide_width - Inches(2.0), Inches(0.35),
                 title, size=14, bold=True, color=SKY)
        _textbox(slide, Inches(1.5), yy + Inches(0.4),
                 prs.slide_width - Inches(2.0), Inches(0.6),
                 body, size=11, color=TEXT)


def slide_outputs(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "7. 산출물 & Quick Start", accent=OK)

    # Outputs tree
    _card(slide, Inches(0.4), Inches(1.2),
          Inches(6.2), Inches(5.6),
          "산출물 (Outputs)",
          [
              "results/YY-MM-DD/",
              "- rx_result.csv  (test_index, ts, channel, length, rx_count)",
              "- logs/HHMMSS_iter<N>_<reason>/",
              "  · bt_test.log / bt_bootstrap.log  (SFTP)",
              "  · FAILURE_SUMMARY.txt  (채널/길이/findings/serial dump)",
              "Plot CSV → min/avg/max + 라인 그래프",
              "환경변수 DTM_RX_RESULT_BASE 로 경로 오버라이드 가능",
          ], accent=OK)

    # Quick start
    _card(slide, Inches(6.8), Inches(1.2),
          Inches(6.1), Inches(5.6),
          "운영자 Quick Start",
          [
              "Dongle USB 연결 → 상단 ● Dongle: COMxx 확인",
              "Channel / Length / Iterations / RX dur / Cooldown 입력",
              "표준 라인: DUT-LINK 모드 그대로 사용",
              "신규 OEM 샘플: MANUAL 토글 → DUT settings… → Select script",
              "단발 측정: START RX → END RX",
              "반복 측정: AUTO RUN (reboot) / AUTO (no reboot)",
              "Open results folder 로 CSV & 실패 로그 확인",
          ], accent=SKY)


def slide_summary(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, prs, "8. 요약 (Summary)", accent=SKY)

    _textbox(slide, Inches(0.6), Inches(1.3),
             prs.slide_width - Inches(1.2), Inches(0.8),
             "Test Sync + Script Selection = 재현 가능한 RX 자동 측정",
             size=24, bold=True, color=SKY, align=PP_ALIGN.CENTER)

    bullets = [
        "Test Sync : Dongle RX 윈도와 DUT TX를 한 시퀀스로 묶어 측정 시점 항상 일치",
        "Script Selection : MANUAL 모드에서 원격 *.sh를 GUI로 선택, OEM 추가 시 코드 변경 X",
        "AUTO Loop : iterations / cooldown 무인 반복 + CSV 자동 누적 + STOP",
        "Failure Capture : rx_count==0 / 알려진 패턴 시 SFTP 로그 + Summary 자동 생성",
        "3 Modes : DUT-LINK / MANUAL / STANDALONE — 라인·디버깅·외부 신호원 모두 대응",
    ]
    _bullets(slide, Inches(1.0), Inches(2.4),
             prs.slide_width - Inches(2.0), Inches(4.5),
             bullets, size=16)

    _textbox(slide, Inches(0.6), Inches(6.7),
             prs.slide_width - Inches(1.2), Inches(0.5),
             "→ 양산 라인·신뢰성 시험·신규 OEM 평가 모두 동일 도구로 커버",
             size=14, bold=True, color=OK, align=PP_ALIGN.CENTER)


# ---------- build ------------------------------------------------------------
def build() -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_why(prs)
    slide_features(prs)
    slide_architecture(prs)
    slide_sequence_single(prs)
    slide_sequence_auto(prs)
    slide_manual_setup(prs)
    slide_outputs(prs)
    slide_summary(prs)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
